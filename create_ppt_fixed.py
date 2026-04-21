from pptx import Presentation

# Create a new blank presentation
prs = Presentation()

slides_content = [
    ("IT Service Management Optimization using Machine Learning",
     "Predicting Ticket Priority and Enhancing IT Incident Response\n\nDimple [Surname]\nB.Tech | DevOps Intern @ Albysol Innovations\nOctober 2025"),

    ("Project Overview",
     "This project applies Machine Learning to improve IT Service Management (ITSM) by predicting ticket priorities and reducing response time."),

    ("Business Problem",
     "• ABC Tech handles 22K–25K IT incidents per year.\n• Existing ITIL processes are mature but slow.\n• Customer survey rated 'Incident Management' as poor.\n• Machine Learning is explored to automate predictions."),

    ("Project Objectives",
     "1. Predict High Priority Tickets (P1/P2)\n2. Forecast Incident Volumes (quarterly/yearly)\n3. Auto-tag Tickets with Correct Priority\n4. Predict Change Request Failures"),

    ("Dataset Description",
     "Dataset of 46,000 records (2012–2014) from ITSM system.\nContains features like impact, urgency, handle time, and reassignment counts."),

    ("Priority Matrix",
     "Business-defined matrix determines ticket priority based on Impact × Urgency levels."),

    ("Data Preprocessing",
     "• Loaded data from CSV\n• Cleaned missing values\n• Encoded categorical columns\n• Selected key features\n• Split data into train/test sets"),

    ("Model Building",
     "Algorithm: Random Forest Classifier\n\nInput Features: Impact, Urgency, No_of_Reassignments, Handle_Time_hrs\nTarget: Priority\n\nModel predicts ticket priority based on input parameters."),

    ("Model Evaluation",
     "Accuracy: ~85–90%\nPrecision: ~0.85\nRecall: ~0.83\n\nConfusion matrix and classification report show strong model performance."),

    ("Flask Web App",
     "Flask-based web app allows live ticket priority prediction via a simple form.\nDeployed locally for testing and demonstration."),

    ("Data Visualization",
     "Generated multiple visualizations:\n• Priority distribution\n• Handle time histogram\n• Impact vs Urgency heatmap\n• Reassignments vs Handle Time scatter plot"),

    ("Forecasting (Optional)",
     "Used ARIMA time series model to forecast future incident volumes for proactive resource planning."),

    ("Key Insights",
     "✅ High urgency + impact = higher priority\n✅ More reassignments → longer handle time\n✅ Database/Network incidents most frequent\n✅ ML improved incident response speed"),

    ("Conclusion",
     "Machine Learning automated ticket prioritization, improved SLA compliance, and enabled predictive analysis.\nFlask web app delivers real-time results."),

    ("Future Enhancements",
     "• Integrate with ServiceNow / Jira\n• Use NLP for ticket description analysis\n• Deploy on Render / AWS / PythonAnywhere\n• Build Streamlit dashboard for visualization")
]

# Add slides with Title & Content layout
for title, content in slides_content:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save properly and close
output_file = "ITSM_Project_Template.pptx"
prs.save(output_file)
print(f"✅ PowerPoint created successfully: {output_file}")
